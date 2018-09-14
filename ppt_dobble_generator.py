import os
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from random import randint
import sys
import comtypes.client

class GenerateDobbleIndex():
    def __init__(self, num_of_pic):
        self.num_of_pic = num_of_pic
    def prime_dobble(self, check = False):
        card_num = self.num_of_pic * (self.num_of_pic-1)+1
        cards = []
        start = 0
        card_list = list(range(start, card_num+start))
        card_list.pop(0)
        key_pic = []
        for pic in range(start, self.num_of_pic+start):
            if pic == start:
                for i in range(self.num_of_pic):
                    select_pic = [pic]
                    for j in range(self.num_of_pic-1):
                        select_pic.append(card_list.pop(0))
                    cards.append(select_pic)
                    if i > 0:
                        key_pic.append(select_pic[1:])
                fixed_pic = key_pic.pop(0)
                #key_pic is the changing part. the elements start from index 0
                #fixed_pic is the elements that stay put
            if pic > start:
                #per card
                for section in range(self.num_of_pic-1):
                    select_pic = [pic, fixed_pic[section]]
                    #inside the card
                    for key in range(self.num_of_pic-2):
                        temp_var = section + (pic-2)*(key+1)
                        while temp_var > self.num_of_pic-2:
                            temp_var = temp_var - self.num_of_pic + 1
                        select_pic.append(key_pic[key][temp_var])
                    cards.append(select_pic)


            card_list = list(range(1, card_num + 1))
            card_list.pop(0)
        if check == True:
            self.check_valid(cards)
        return cards

    def non_prime_dobble(self, check = False):
        card_num = self.num_of_pic * (self.num_of_pic - 1) + 1
        cards = []
        start = 0
        card_list = list(range(start, card_num + start))
        card_list.pop(0)
        key_list = [[[5, 9, 13, 17], [6, 10, 15, 20], [7, 11, 16, 18], [8, 12, 14, 19]], [[5, 10, 14, 18], [6, 9, 16, 19], [7, 12, 15, 17], [8, 11, 13, 20]], [[5, 11, 15, 19], [6, 12, 13, 18], [7, 9, 14, 20], [8, 10, 16, 17]], [[5, 12, 16, 20], [6, 11, 14, 17], [7, 10, 13, 19], [8, 9, 15, 18]]]
        for pic in range(start, self.num_of_pic + start):
            if pic == start:
                for i in range(self.num_of_pic):
                    select_pic = [pic]
                    for j in range(self.num_of_pic - 1):
                        select_pic.append(card_list.pop(0))
                    cards.append(select_pic)
            else:
                for i in range(self.num_of_pic-1):
                    select_pic = [pic] + key_list[pic-2][i]
                    cards.append(select_pic)
        if check == True:
            self.check_valid(cards)
        return cards


    def check_valid(self, c):
        valid = True
        count = 0
        for i in c:
            for j in c:
                dup = len(set(i) & set(j))
                if 1< dup < len(c[0]):
                    valid = False
                    count += 1
        if valid == False:
            print(len(c[0]))
            print(count//2, "쌍의 개수가 도블 조건을 만족하지 않습니다.")

    def make_dobble_set(self):
        if self.num_of_pic == 5:
            return self.non_prime_dobble(self.num_of_pic)
        else:
            return self.prime_dobble(self.num_of_pic)

    def get_dobble_index(self):
        cards = self.make_dobble_set()
        cards_index = [pic for card in cards for pic in card]
        return cards_index

class PptCardMaker():
    def __init__(self, words_pics, num):
        self.words_pics = words_pics
        self.num = num
        try:
            self.desktop = os.path.join(os.path.join(os.environ['USERPROFILE']), 'desktop')
            os.listdir(self.desktop)
        except:
            self.desktop = '.'
        self.words = [word[0] for word in self.words_pics]
        self.pics = [pic[1] for pic in self.words_pics]

    def make_card_with_picture(self):
        prs = Presentation('template_for_dobble_cards.pptx')
        # prs.slide_width = 11887200
        # prs.slide_height = 6686550
        for i in range((len(self.pics)-1)//self.num+1):
            card_layout = prs.slide_layouts[self.num-3]
            slide = prs.slides.add_slide(card_layout)

            for j, shape in enumerate(slide.placeholders):
                pic = shape.insert_picture(self.pics[self.num*i + j])
                pic.crop_left = 0
                pic.crop_right = 0
                pic.crop_bottom = 0
                pic.crop_top = 0

        file_path = os.path.join(self.desktop, 'dobble_{}cards_picture.pptx'.format(self.num))
        prs.save(file_path)
        return [file_path]

    def make_card_with_word(self):
        prs = Presentation('template_for_dobble_cards.pptx')
        for i in range((len(self.pics)-1)//self.num+1):
            card_layout = prs.slide_layouts[self.num-3]
            slide = prs.slides.add_slide(card_layout)

            for j, shape in enumerate(slide.placeholders):
                text_frame = shape.text_frame
                text_frame.text = self.words[self.num * i + j]
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        file_path = os.path.join(self.desktop, 'dobble_{}cards_word.pptx'.format(self.num))
        prs.save(file_path)
        return [file_path]



class ConvertPptToPng():
    def __init__(self, input_file):
        self.input_file = input_file
    def init_powerpoint(self):
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        return powerpoint

    def ppt_to_png(self, powerpoint, inputFileName, outputFileName, formatType = 18):
        deck = powerpoint.Presentations.Open(inputFileName, WithWindow = False)
        deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
        deck.Close()

    def convert(self):
        comtypes.CoInitialize()
        powerpoint = self.init_powerpoint()
        self.ppt_to_png(powerpoint, self.input_file, self.input_file)
        powerpoint.Quit()
        comtypes.CoUninitialize()


import threading
if __name__ == '__main__':
    con = ConvertPptToPng("C:\\Users\\hoetaekpro\\Desktop\\dobble_3cards_word.pptx")
    p = threading.Thread(target=con.convert)
    p.start()
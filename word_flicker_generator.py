from pptx import Presentation
from PIL import ImageFont
from PIL import Image
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Length, Cm
import requests
from bs4 import BeautifulSoup
import re
import os
import re
from syllable import get_syllable_divided

class PptWordFlickerMaker():
    def __init__(self, words_pics, num_slide, path):
        self.words_pics = words_pics
        self.num_slide = num_slide
        self.words = [word[0] for word in self.words_pics]
        self.pics = [pic[1] for pic in self.words_pics]
        self.path = path

    def make_word_flicker_slide(self):
        prs = Presentation('template_for_word_flicker.pptx')
        for num in self.num_slide:
            if num != 0 or num != 1:
                prs.slides.add_slide(prs.slide_layouts[-2])
                prs.slides.add_slide(prs.slide_layouts[-1])
            card_layout = prs.slide_layouts[num]
            for word, pic in zip(self.words, self.pics):

                slide = prs.slides.add_slide(card_layout)
                for shape in slide.shapes:
                    if str(shape.placeholder_format.type) == 'BODY (2)':
                        text_frame = shape.text_frame
                        if num == 1:
                            if shape.placeholder_format.idx == 10:
                                p = text_frame.paragraphs[0]
                                run = p.add_run()
                                run.text = word
                            else:
                                p = text_frame.paragraphs[0]
                                run = p.add_run()

                                req = requests.get('http://endic.naver.com/small_search.nhn?query=' + word)
                                html = req.text
                                soup = BeautifulSoup(html, 'html.parser')
                                meanings = soup.select('span.fnt_k05')
                                text = meanings[0].text
                                parenthesis = re.compile(r'(\s)?\(.*\)(\s)?')
                                bracket = re.compile(r'(\s)?\[.*\](\s)?')
                                text = re.sub(parenthesis, '', text)
                                text = re.sub(bracket, '', text)
                                run.text = text
                        elif num == 4:
                            p = text_frame.paragraphs[0]
                            run = p.add_run()

                            req = requests.get('http://endic.naver.com/small_search.nhn?query=' + word)
                            html = req.text
                            soup = BeautifulSoup(html, 'html.parser')
                            meanings = soup.select('span.fnt_k05')
                            text = meanings[0].text
                            parenthesis = re.compile(r'(\s)?\(.*\)(\s)?')
                            bracket = re.compile(r'(\s)?\[.*\](\s)?')
                            text = re.sub(parenthesis, '', text)
                            text = re.sub(bracket, '', text)
                            run.text = text
                        elif num == 5:
                            syllable_divided = ' '.join(get_syllable_divided(w) for w in word.split())
                            if shape.placeholder_format.idx == 10:
                                p = text_frame.paragraphs[0]
                                run = p.add_run()
                                run.text = word
                            elif shape.placeholder_format.idx == 11:
                                p = text_frame.paragraphs[0]
                                run = p.add_run()
                                run.text = syllable_divided
                            elif shape.placeholder_format.idx == 12:
                                p = text_frame.paragraphs[0]
                                run = p.add_run()
                                syllable_num = len(re.split(r'[ -]', syllable_divided))
                                if syllable_num == 1:
                                    run.text = str(syllable_num) + ' syllable'
                                else:
                                    run.text = str(syllable_num) + ' syllables'
                        else:
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = word


                for shape in slide.shapes:
                    if str(shape.placeholder_format.type) == 'PICTURE (18)':
                        # im = Image.open('whatever.png')
                        # width, height = im.size
                        picture = shape.insert_picture(pic)
                        # add_picture(image_file, left, top, width=None, height=No)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        picture.crop_bottom = 0
                        picture.crop_top = 0


        file_path = os.path.join(self.path, '단어깜빡이.pptx')
        file_path = os.path.abspath(file_path)
        prs.save(file_path)
        return file_path

if __name__=='__main__':
    ppt = PptWordFlickerMaker([['brush', 'C:\\Users\\user\\Desktop\\구글이미지\\brush\\6. bad-breath-brush.jpg'], ['exercise', 'C:\\Users\\user\\Desktop\\구글이미지\\exercise\\7. well_howtostartrunning_promo-largehorizontaljumbo.jpg'],
                                ['healthy', 'C:\\Users\\user\\Desktop\\구글이미지\\healthy\\3. workout-composition-with-healthy-food_23-2147692092.jpg'], ['often', 'C:\\Users\\user\\Desktop\\구글이미지\\often\\5. often.jpg'], ['tooth', 'C:\\Users\\user\\Desktop\\구글이미지\\tooth\\3. istocktooth.jpg'], ['teeth', 'C:\\Users\\user\\Desktop\\구글이미지\\teeth\\7. 011316_3dteeth_thumb_large.jpg'], ['twice', 'C:/Users/user/Desktop/구글이미지/x2/x2.png'],
                                ['three times', 'C:\\Users\\user\\Desktop\\구글이미지\\threetimes\\x3.png'], ['breakfast', 'C:\\Users\\user\\Desktop\\구글이미지\\breakfast\\5. 7562ab71-9093-41e3-9534-6509501370ad--2018-0309_wholeflour-breakfast-cookie_3x2_rocky-luten_033.jpg'], ['every week', 'C:\\Users\\user\\Desktop\\구글이미지\\everyweek\\8. small-scale-sabbaticals.jpg']], [0, 2, 3], '.')
    ppt.make_word_flicker_slide()
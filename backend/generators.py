import math
import wave
from base64 import b64decode
from collections.abc import Sequence
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import httpx
from docx import Document
from docx.document import Document as DocxDocument
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.oxml.xmlchemy import BaseOxmlElement, OxmlElement
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Inches
from pptx.util import Pt as PptPt

from backend.image_search import IMAGE_REQUEST_HEADERS
from backend.schemas import (
    DobbleRequest,
    FlickerRequest,
    WordSearchRequest,
    WorksheetRequest,
)

FLICKER_AUTO_ADVANCE_MS = 1000
FLICKER_TRANSITION_DURATION_MS = 700
FLICKER_TRANSITION_SPEED = "med"
FLICKER_REVEAL_DELAY_MS = 500
FLICKER_REVEAL_DURATION_MS = 500
P14_NS_URI = "http://schemas.microsoft.com/office/powerpoint/2010/main"
WAV_CONTENT_TYPE = "audio/x-wav"
WORD_SEARCH_STUDENT_INSTRUCTION = "퀴즈 힌트를 풀고 낱말을 찾아 동그라미 하세요."
WORKSHEET_STUDENT_INSTRUCTION = "그림을 보고 단어를 읽은 뒤 빈칸에 따라 쓰세요."


async def make_flicker_pptx(request: FlickerRequest) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    transition_audio = make_transition_audio_part(presentation)

    for item in request.items:
        for template in request.templates:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            sound_rel_id = slide.part.relate_to(transition_audio, RT.AUDIO)
            add_auto_advance_transition(slide, sound_rel_id)
            add_slide_frame(slide, item.word)
            word_shape: BaseShape | None = None
            image_shape: BaseShape | None = None

            if template in {"word", "word-image"}:
                word_shape = add_center_text(
                    slide,
                    item.word,
                    top=0.75 if template == "word-image" else 2.45,
                )

            if template in {"image", "word-image"}:
                image = await resolve_image(item.image)
                if image:
                    image_shape = add_image(
                        slide,
                        image,
                        left=4.65,
                        top=2.15,
                        width=4.0,
                        height=3.35,
                    )
                else:
                    add_placeholder(slide, "사진 없음")

            if template == "word-image" and word_shape is not None and image_shape is not None:
                add_word_image_reveal_timing(slide, word_shape.shape_id, image_shape.shape_id)

            if template == "blank":
                add_center_text(slide, "", top=2.45)

    return save_presentation(presentation)


async def make_dobble_pptx(request: DobbleRequest) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(10)

    for index, card in enumerate(request.cards):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_dobble_card_frame(slide)
        add_small_label(slide, f"도블 카드 {index + 1}")
        symbol_layout = card_symbol_layout(len(card), index)
        symbol_size = card_symbol_size(len(card))
        for item, (left, top, rotation) in zip(card, symbol_layout, strict=False):
            image = await resolve_image(item.image) if request.display_mode != "word" else None
            if image and request.display_mode != "word":
                image_shape = add_image(
                    slide,
                    image,
                    left=left,
                    top=top,
                    width=symbol_size,
                    height=symbol_size,
                )
                image_shape.rotation = rotation
            elif request.display_mode != "word":
                symbol = add_dobble_initial_symbol(
                    slide,
                    item.word,
                    left=left,
                    top=top,
                    size=symbol_size,
                )
                symbol.rotation = rotation
            if request.display_mode == "image":
                continue

            text_left = left - 0.3 if request.display_mode == "image-word" else left - 0.45
            text_top = (
                top + symbol_size + 0.03
                if request.display_mode == "image-word"
                else top + symbol_size * 0.3
            )
            text_width = (
                symbol_size + 0.6 if request.display_mode == "image-word" else symbol_size + 0.9
            )
            text_height = 0.52 if request.display_mode == "image-word" else symbol_size * 0.55
            textbox = slide.shapes.add_textbox(
                Inches(text_left),
                Inches(text_top),
                Inches(text_width),
                Inches(text_height),
            )
            textbox.rotation = rotation
            paragraph = textbox.text_frame.paragraphs[0]
            paragraph.text = item.word
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = PptPt(11 if request.display_mode == "image-word" else 18)
            paragraph.font.bold = request.display_mode == "word"

    return save_presentation(presentation)


async def make_worksheet_docx(request: WorksheetRequest) -> bytes:
    document = Document()
    section = document.sections[0]
    section.top_margin = Cm(1)
    section.bottom_margin = Cm(0.8)
    section.left_margin = Cm(1.3)
    section.right_margin = Cm(1.3)

    heading = document.add_heading("단어 활동지", 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_student_info(document, request.grade, request.class_number)
    instruction = document.add_paragraph(WORKSHEET_STUDENT_INSTRUCTION)
    instruction.alignment = WD_ALIGN_PARAGRAPH.CENTER
    instruction.runs[0].bold = True

    rows = (len(request.items) + request.columns - 1) // request.columns
    table = document.add_table(rows=rows, cols=request.columns, style="Table Grid")
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            item_index = row_index * request.columns + col_index
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            paragraph = cell.paragraphs[0]
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if item_index >= len(request.items):
                continue

            item = request.items[item_index]
            image = await resolve_image(item.image)
            if image:
                paragraph.add_run().add_picture(image, width=Cm(3.0))
            else:
                paragraph.add_run("사진 없음")

            text = split_syllables(item.word) if request.syllables else item.word
            word_paragraph = cell.add_paragraph()
            word_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            word_run = word_paragraph.add_run(text)
            word_run.font.name = "Arial"
            word_run.font.size = Pt(15)
            word_run.bold = True

            practice_label = cell.add_paragraph("써 보기")
            practice_label.alignment = WD_ALIGN_PARAGRAPH.LEFT
            practice_label.runs[0].font.size = Pt(9)
            for _ in range(2):
                line_paragraph = cell.add_paragraph()
                line_run = line_paragraph.add_run(" " * 24)
                line_run.underline = True

    return save_document(document)


async def make_word_search_docx(request: WordSearchRequest) -> bytes:
    document = Document()
    section = document.sections[0]
    section.top_margin = Cm(1)
    section.bottom_margin = Cm(0.8)
    section.left_margin = Cm(1.7)
    section.right_margin = Cm(1.7)

    heading = document.add_heading(request.title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_student_info(document, request.grade, request.class_number)
    instruction = document.add_paragraph(WORD_SEARCH_STUDENT_INSTRUCTION)
    instruction.alignment = WD_ALIGN_PARAGRAPH.CENTER
    instruction.runs[0].bold = True

    width = max(len(request.grid[0]), 1)
    table = document.add_table(rows=len(request.grid), cols=width, style="Table Grid")
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row_index, row in enumerate(request.grid):
        for col_index, letter in enumerate(row):
            cell = table.rows[row_index].cells[col_index]
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            paragraph = cell.paragraphs[0]
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = paragraph.add_run(letter)
            run.bold = True

    if request.hints:
        document.add_paragraph("")
        hint_heading = document.add_paragraph("퀴즈 힌트")
        hint_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        hint_heading.runs[0].bold = True

        hint_table = document.add_table(rows=len(request.hints) + 1, cols=3, style="Table Grid")
        hint_table.alignment = WD_TABLE_ALIGNMENT.CENTER

        header_cells = hint_table.rows[0].cells
        for header_cell, label in zip(
            header_cells, ["번호", "퀴즈 힌트", "찾은 낱말"], strict=True
        ):
            header_cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            header_paragraph = header_cell.paragraphs[0]
            header_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            header_run = header_paragraph.add_run(label)
            header_run.font.size = Pt(10)
            header_run.bold = True

        for index, item in enumerate(request.hints, start=1):
            cells = hint_table.rows[index].cells
            for cell in cells:
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

            number_paragraph = cells[0].paragraphs[0]
            number_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            number_run = number_paragraph.add_run(str(index))
            number_run.font.size = Pt(10)
            number_run.bold = True

            clue_paragraph = cells[1].paragraphs[0]
            clue_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            clue_run = clue_paragraph.add_run((item.clue or "").strip())
            clue_run.font.size = Pt(11)

            answer_paragraph = cells[2].paragraphs[0]
            answer_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            answer_run = answer_paragraph.add_run(" " * 18)
            answer_run.underline = True

    return save_document(document)


async def make_zip(files: Sequence[tuple[str, bytes]]) -> bytes:
    buffer = BytesIO()
    with ZipFile(buffer, "w", ZIP_DEFLATED) as archive:
        for name, content in files:
            archive.writestr(name, content)
    return buffer.getvalue()


async def resolve_image(value: str | None) -> BytesIO | None:
    if not value:
        return None

    if value.startswith("data:image/"):
        _, encoded = value.split(",", 1)
        return normalize_image(BytesIO(b64decode(encoded)))

    if value.startswith(("http://", "https://")):
        try:
            async with httpx.AsyncClient(
                timeout=15,
                headers=IMAGE_REQUEST_HEADERS,
                follow_redirects=True,
            ) as client:
                response = await client.get(value)
                response.raise_for_status()
                content_type = response.headers.get("content-type", "")
                if not content_type.startswith("image/"):
                    return None
                return normalize_image(BytesIO(response.content))
        except (httpx.HTTPError, UnidentifiedImageError, OSError):
            return None

    return None


def normalize_image(image: BytesIO, size: tuple[int, int] = (800, 600)) -> BytesIO:
    image.seek(0)
    with Image.open(image) as source:
        source.thumbnail(size, Image.Resampling.LANCZOS)
        frame = Image.new("RGB", size, "#ffffff")

        if source.mode in {"RGBA", "LA"} or (source.mode == "P" and "transparency" in source.info):
            normalized = source.convert("RGBA")
            position = ((size[0] - normalized.width) // 2, (size[1] - normalized.height) // 2)
            frame.paste(normalized, position, normalized)
        else:
            normalized = source.convert("RGB")
            position = ((size[0] - normalized.width) // 2, (size[1] - normalized.height) // 2)
            frame.paste(normalized, position)

    buffer = BytesIO()
    frame.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


def add_slide_frame(slide: Slide, label: str) -> None:
    add_small_label(slide, label)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.75),
        Inches(6.8),
        Inches(12.55),
        Inches(6.8),
    )
    line.line.color.rgb = RGBColor(235, 235, 235)  # type: ignore[no-untyped-call]


def make_transition_audio_part(presentation: PptxPresentation) -> Part:
    package = presentation.part.package
    return Part(
        package.next_partname("/ppt/media/audio%d.wav"),
        WAV_CONTENT_TYPE,
        package,
        blob=build_transition_sound(),
    )


def build_transition_sound() -> bytes:
    sample_rate = 11025
    duration_seconds = 0.16
    frame_count = int(sample_rate * duration_seconds)
    frames = bytearray()

    for index in range(frame_count):
        progress = index / frame_count
        envelope = (1 - progress) ** 2
        frequency = 1200 - (420 * progress)
        sample = 128 + int(
            72 * envelope * math.sin(2 * math.pi * frequency * (index / sample_rate))
        )
        frames.append(min(255, max(0, sample)))

    buffer = BytesIO()
    with wave.open(buffer, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(1)
        wav_file.setframerate(sample_rate)
        wav_file.writeframes(bytes(frames))

    return buffer.getvalue()


def add_auto_advance_transition(slide: Slide, sound_rel_id: str) -> None:
    transition = OxmlElement("p:transition")
    transition.set("spd", FLICKER_TRANSITION_SPEED)
    transition.set(f"{{{P14_NS_URI}}}dur", str(FLICKER_TRANSITION_DURATION_MS))
    transition.set("advClick", "0")
    transition.set("advTm", str(FLICKER_AUTO_ADVANCE_MS))
    transition.append(OxmlElement("p:fade"))
    transition.append(make_transition_sound_action(sound_rel_id))

    slide_element = slide.element
    for existing_transition in slide_element.findall(qn("p:transition")):
        slide_element.remove(existing_transition)

    insert_index = len(slide_element)
    for index, child in enumerate(slide_element):
        if child.tag in {qn("p:timing"), qn("p:extLst")}:
            insert_index = index
            break

    slide_element.insert(insert_index, transition)


def make_transition_sound_action(sound_rel_id: str) -> BaseOxmlElement:
    sound_action = OxmlElement("p:sndAc")
    start_sound = OxmlElement("p:stSnd")
    sound = OxmlElement("p:snd")
    sound.set(qn("r:embed"), sound_rel_id)
    sound.set("name", "whoosh.wav")
    start_sound.append(sound)
    sound_action.append(start_sound)
    return sound_action


def add_word_image_reveal_timing(
    slide: Slide,
    word_shape_id: int,
    image_shape_id: int,
) -> None:
    timing = parse_xml(
        f"""
        <p:timing {nsdecls("p")}>
          <p:tnLst>
            <p:par>
              <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                <p:childTnLst>
                  <p:seq concurrent="1" nextAc="seek">
                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                      <p:childTnLst>
                        <p:par>
                          <p:cTn id="3" fill="hold">
                            <p:stCondLst>
                              <p:cond delay="{FLICKER_REVEAL_DELAY_MS}"/>
                            </p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="4" presetID="10" presetClass="exit"
                                  presetSubtype="0" fill="hold" grpId="0"
                                  nodeType="afterEffect">
                                  <p:stCondLst>
                                    <p:cond delay="0"/>
                                  </p:stCondLst>
                                  <p:childTnLst>
                                    <p:animEffect transition="out" filter="fade">
                                      <p:cBhvr>
                                        <p:cTn id="5" dur="{FLICKER_REVEAL_DURATION_MS}"/>
                                        <p:tgtEl>
                                          <p:spTgt spid="{word_shape_id}"/>
                                        </p:tgtEl>
                                      </p:cBhvr>
                                    </p:animEffect>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                              <p:par>
                                <p:cTn id="6" presetID="10" presetClass="entr"
                                  presetSubtype="0" fill="hold" grpId="0"
                                  nodeType="withEffect">
                                  <p:stCondLst>
                                    <p:cond delay="0"/>
                                  </p:stCondLst>
                                  <p:childTnLst>
                                    <p:animEffect transition="in" filter="fade">
                                      <p:cBhvr>
                                        <p:cTn id="7" dur="{FLICKER_REVEAL_DURATION_MS}"/>
                                        <p:tgtEl>
                                          <p:spTgt spid="{image_shape_id}"/>
                                        </p:tgtEl>
                                      </p:cBhvr>
                                    </p:animEffect>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                      </p:childTnLst>
                    </p:cTn>
                    <p:prevCondLst>
                      <p:cond evt="onPrev" delay="0">
                        <p:tgtEl>
                          <p:sldTgt/>
                        </p:tgtEl>
                      </p:cond>
                    </p:prevCondLst>
                    <p:nextCondLst>
                      <p:cond evt="onNext" delay="0">
                        <p:tgtEl>
                          <p:sldTgt/>
                        </p:tgtEl>
                      </p:cond>
                    </p:nextCondLst>
                  </p:seq>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:tnLst>
          <p:bldLst>
            <p:bldP spid="{word_shape_id}" grpId="0" build="p"/>
            <p:bldP spid="{image_shape_id}" grpId="0" build="p"/>
          </p:bldLst>
        </p:timing>
        """
    )

    slide_element = slide.element
    for existing_timing in slide_element.findall(qn("p:timing")):
        slide_element.remove(existing_timing)

    insert_index = len(slide_element)
    for index, child in enumerate(slide_element):
        if child.tag == qn("p:extLst"):
            insert_index = index
            break

    slide_element.insert(insert_index, timing)


def add_small_label(slide: Slide, label: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(4), Inches(0.35))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = label
    paragraph.font.size = PptPt(11)
    paragraph.font.bold = True


def add_center_text(slide: Slide, text: str, top: float) -> BaseShape:
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(11.3), Inches(1.2))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = PptPt(44)
    paragraph.font.bold = True
    return textbox


def add_placeholder(slide: Slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(4.65), Inches(3.25), Inches(4), Inches(0.6))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = PptPt(18)


def add_image(
    slide: Slide,
    image: BytesIO,
    left: float,
    top: float,
    width: float,
    height: float,
) -> BaseShape:
    image.seek(0)
    return slide.shapes.add_picture(
        image,
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )


def add_dobble_card_frame(slide: Slide) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.55),
        Inches(0.55),
        Inches(8.9),
        Inches(8.9),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)  # type: ignore[no-untyped-call]
    card.line.color.rgb = RGBColor(224, 224, 220)  # type: ignore[no-untyped-call]
    card.line.width = PptPt(2)


def add_dobble_initial_symbol(
    slide: Slide,
    word: str,
    left: float,
    top: float,
    size: float,
) -> BaseShape:
    symbol = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(size),
        Inches(size),
    )
    symbol.fill.solid()
    symbol.fill.fore_color.rgb = RGBColor(240, 241, 242)  # type: ignore[no-untyped-call]
    symbol.line.color.rgb = RGBColor(255, 255, 255)  # type: ignore[no-untyped-call]
    symbol.line.width = PptPt(2)

    text_frame = symbol.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.text = word.strip()[0].upper()
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = PptPt(30)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(23, 23, 23)  # type: ignore[no-untyped-call]
    return symbol


def card_positions(count: int) -> list[tuple[float, float]]:
    return [(left, top) for left, top, _rotation in card_symbol_layout(count)]


def card_symbol_layout(count: int, card_index: int = 0) -> list[tuple[float, float, int]]:
    size = card_symbol_size(count)
    center_x = 5.0
    center_y = 5.08
    fixed_centers: dict[int, list[tuple[float, float]]] = {
        3: [(5.0, 2.68), (7.16, 6.33), (2.84, 6.33)],
        4: [(3.14, 3.21), (6.86, 3.21), (6.86, 6.95), (3.14, 6.95)],
        5: [(5.0, 5.08), (2.84, 2.92), (7.16, 2.92), (7.16, 7.24), (2.84, 7.24)],
    }
    fixed_layout = fixed_centers.get(count)

    if fixed_layout is not None:
        return [
            (
                left - size / 2,
                top - size / 2,
                dobble_symbol_rotation(left, top, index, card_index),
            )
            for index, (left, top) in enumerate(fixed_layout)
        ]

    radius = 3.16
    positions: list[tuple[float, float, int]] = []

    for index in range(count):
        angle = -math.pi / 2 + (index * 2 * math.pi) / count
        symbol_center_x = center_x + math.cos(angle) * radius
        symbol_center_y = center_y + math.sin(angle) * radius
        left = symbol_center_x - size / 2
        top = symbol_center_y - size / 2
        positions.append(
            (
                left,
                top,
                dobble_symbol_rotation(symbol_center_x, symbol_center_y, index, card_index),
            )
        )

    return positions


def dobble_symbol_rotation(x: float, y: float, index: int, card_index: int) -> int:
    dx = x - 5.0
    dy = y - 5.08
    jitter = dobble_rotation_jitter(index, card_index)

    if math.hypot(dx, dy) < 0.01:
        return jitter % 360

    return round((math.degrees(math.atan2(dy, dx)) - 90 + jitter) % 360)


def dobble_rotation_jitter(index: int, card_index: int) -> int:
    jitters = [-7, 6, -4, 8, -6, 5, 3]
    return jitters[(card_index * 3 + index * 2) % len(jitters)]


def card_symbol_size(count: int) -> float:
    if count == 3:
        return 2.65
    if count == 4:
        return 2.32
    if count == 5:
        return 2.02
    if count >= 8:
        return 1.32
    if count >= 6:
        return 1.45
    return 1.62


def split_syllables(word: str) -> str:
    return " · ".join(word)


def add_student_info(document: DocxDocument, grade: int, class_number: int) -> None:
    table = document.add_table(rows=1, cols=3, style="Table Grid")
    table.alignment = WD_TABLE_ALIGNMENT.RIGHT
    table.autofit = False

    fields = [
        ("학년", str(grade), Cm(2.0)),
        ("반", str(class_number), Cm(2.0)),
        ("이름", "", Cm(4.6)),
    ]
    for cell, (label, value, width) in zip(table.rows[0].cells, fields, strict=True):
        cell.width = width
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        paragraph = cell.paragraphs[0]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        label_run = paragraph.add_run(label)
        label_run.font.size = Pt(9)
        label_run.bold = True
        if value:
            value_run = paragraph.add_run(f" {value}")
            value_run.font.size = Pt(14)
            value_run.bold = True

    document.add_paragraph("")


def save_presentation(presentation: PptxPresentation) -> bytes:
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def save_document(document: DocxDocument) -> bytes:
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()
